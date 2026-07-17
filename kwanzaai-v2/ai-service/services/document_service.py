import fitz
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import csv
import re
import json

class DocumentService:
    #chunk teks
    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 100):
        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk.strip())
            start += chunk_size - overlap
        return chunks

    #clean teks
    def normalize_text(self, text: str):
        if not text:
            return ""
        text = text.replace("\t", " ")
        text = text.replace("\r", "\n")
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r"[ ]{2,}", " ", text)
        text = re.sub(r"\u00A0", " ", text)
        return text.strip()

    def clean_text(self, text: str):
        if not text:
            return ""
        text = re.sub(r"[\x00-\x08\x0B-\x1F\x7F]", "", text)
        text = re.sub(r"[•●▪►■◆]", "-", text)
        text = re.sub(r"[“”]", "\"", text)
        text = re.sub(r"[‘’]", "'", text)
        text = re.sub(r"…", "...", text)
        text = re.sub(r"­", "", text)
        text = re.sub(r"\n +", "\n", text)
        text = re.sub(r" +\n", "\n", text)
        return text.strip()

    def fix_structure(self, text: str):
        lines = []
        for line in text.split("\n"):
            line = line.strip()
            if line == "":
                continue
            lines.append(line)
        return "\n".join(lines)

    def final_clean(self, text: str):
        text = re.sub(r"\n{2,}", "\n", text)
        text = re.sub(r" {2,}", " ", text)
        return text.strip()

    #extrak teks dari dokument
    def extract_text(self, file_path: str, content_type: str) -> str:
        if content_type == "pdf":
            return self.extract_pdf(file_path)
        elif content_type == "docx":
            return self.extract_docx(file_path)
        elif content_type == "pptx":
            return self.extract_pptx(file_path)
        elif content_type == "xlsx":
            return self.extract_xlsx(file_path)
        elif content_type == "txt":
            return self.extract_txt(file_path)
        elif content_type == "csv":
            return self.extract_csv(file_path)
        elif content_type == "json":
            return self.extract_json(file_path)
        elif content_type == "html":
            return self.extract_html(file_path)
        elif content_type == "md":
            return self.extract_txt(file_path)

    def extract_pdf(self, file_path: str):
        text = ""
        pdf = fitz.open(file_path)
        for page in pdf:
            text += page.get_text()
        pdf.close()
        return text

    def extract_docx(self, file_path: str):
        document = Document(file_path)
        text = []
        for paragraph in document.paragraphs:
            text.append(paragraph.text)
        return "\n".join(text)

    def extract_pptx(self, file_path: str):
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def extract_xlsx(self, file_path: str):
        workbook = load_workbook(file_path)
        text = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                values = []
                for cell in row:
                    if cell is not None:
                        values.append(str(cell))
                text.append(" | ".join(values))
        return "\n".join(text)

    def extract_txt(self, file_path: str):
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()

    def extract_csv(self, file_path: str):
        text = []
        with open(file_path, newline="", encoding="utf-8") as file:
            reader = csv.reader(file)
            for row in reader:
                text.append(" | ".join(row))
        return "\n".join(text)

    def extract_json(self, file_path: str):
        with open(file_path, "r", encoding="utf-8") as file:
            data = json.load(file)
        return json.dumps(
            data,
            indent=2,
            ensure_ascii=False
        )

    def extract_html(self, file_path: str):
        with open(file_path, "r", encoding="utf-8") as file:
            html = file.read()
        soup = BeautifulSoup(
            html,
            "html.parser"
        )
        return soup.get_text(separator="\n")